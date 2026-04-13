"""
School portal report exports: PDF (ReportLab), Excel (openpyxl), PowerPoint (python-pptx).
Used by /api/school-portal/export/* (Bearer school session).
"""
from __future__ import annotations

import logging
from datetime import datetime, timedelta, timezone
from io import BytesIO
from typing import Any, Dict, List

from database.external_db import make_supabase_request
from services.school_portal_scope import (
    portal_network_scope,
    fetch_users_registration_by_school_pks,
    fetch_school_payments_by_codes,
)

logger = logging.getLogger(__name__)

REVENUE_PER_STUDENT = 10.0
SCHOOL_SHARE_PERCENT = 0.30
NERDX_SHARE_PERCENT = 0.70


def _now_utc():
    return datetime.now(timezone.utc)


def _parse_utc(value):
    if value is None:
        return None
    try:
        s = str(value).strip()
        if s.endswith("Z"):
            s = s[:-1] + "+00:00"
        dt = datetime.fromisoformat(s)
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)
        return dt.astimezone(timezone.utc)
    except Exception:
        return None


def _load_portal_students(school_row: Dict[str, Any]) -> List[Dict[str, Any]]:
    db_ids, _ = portal_network_scope(school_row)
    return fetch_users_registration_by_school_pks(
        db_ids,
        "id,name,surname,student_level,subscription_expires_at,xp,streak,last_daily_reset",
    )


def build_overview_snapshot(school_row: Dict[str, Any]) -> Dict[str, Any]:
    """Lightweight overview numbers for PDF/PPTX/AI."""
    _, pay_codes = portal_network_scope(school_row)
    students = _load_portal_students(school_row)
    total = len(students)
    now = _now_utc()
    thirty = now - timedelta(days=30)
    active = sum(
        1
        for st in students
        if _parse_utc(st.get("last_daily_reset")) and _parse_utc(st.get("last_daily_reset")) >= thirty
    )
    sub_exp = _parse_utc(school_row.get("subscription_expires_at"))
    sub_ok = bool(sub_exp and sub_exp >= now)
    days_left = max(0, (sub_exp - now).days) if sub_exp and sub_exp >= now else 0
    payments = fetch_school_payments_by_codes(pay_codes, "amount,status")
    total_paid = sum(float(p.get("amount", 0)) for p in payments if p.get("status") in ("paid", "verified"))
    rev_monthly = total * REVENUE_PER_STUDENT
    school_share = rev_monthly * SCHOOL_SHARE_PERCENT
    nerdx_share = rev_monthly * NERDX_SHARE_PERCENT
    return {
        "total_students": total,
        "active_30d": active,
        "subscription_active": sub_ok,
        "days_remaining": days_left,
        "revenue_monthly_model": rev_monthly,
        "school_share_model": school_share,
        "nerdx_share_model": nerdx_share,
        "total_paid": total_paid,
        "amount_due": max(0, nerdx_share - total_paid),
    }


def generate_summary_pdf(school_row: Dict[str, Any]) -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    snap = build_overview_snapshot(school_row)
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    w, h = letter
    y = h - 50
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, y, "NerdX — School summary report")
    y -= 28
    c.setFont("Helvetica", 11)
    c.drawString(50, y, f"School: {school_row.get('name') or ''} ({school_row.get('school_id')})")
    y -= 16
    c.drawString(50, y, f"Generated: {_now_utc().strftime('%Y-%m-%d %H:%M UTC')}")
    y -= 28
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y, "Portal learners")
    y -= 18
    c.setFont("Helvetica", 11)
    for line in [
        f"Total students: {snap['total_students']}",
        f"Active (30 days): {snap['active_30d']}",
        f"Subscription active: {'Yes' if snap['subscription_active'] else 'No'}",
        f"Days remaining: {snap['days_remaining']}",
        f"Model monthly revenue (USD): ${snap['revenue_monthly_model']:.2f}",
        f"Model school share (30%): ${snap['school_share_model']:.2f}",
        f"Model platform fee (70%): ${snap['nerdx_share_model']:.2f}",
        f"Verified payments (USD): ${snap['total_paid']:.2f}",
        f"Amount due (model vs paid): ${snap['amount_due']:.2f}",
    ]:
        c.drawString(50, y, line)
        y -= 16
        if y < 80:
            c.showPage()
            y = h - 50
            c.setFont("Helvetica", 11)
    c.setFont("Helvetica-Oblique", 9)
    c.drawString(50, 50, "Figures reflect NerdX portal data and the standard per-learner revenue model.")
    c.save()
    return buf.getvalue()


def generate_students_xlsx(school_row: Dict[str, Any]) -> bytes:
    try:
        from openpyxl import Workbook
    except ImportError as e:
        logger.error("openpyxl missing: %s", e)
        raise RuntimeError("Excel export unavailable (openpyxl not installed)") from e

    students = _load_portal_students(school_row)
    wb = Workbook()
    ws = wb.active
    ws.title = "Students"
    headers = ["ID", "First name", "Surname", "Level", "XP", "Streak", "Subscription expires", "Last active"]
    ws.append(headers)
    for st in students:
        ws.append(
            [
                st.get("id"),
                st.get("name") or "",
                st.get("surname") or "",
                st.get("student_level") or "",
                st.get("xp") or 0,
                st.get("streak") or 0,
                str(st.get("subscription_expires_at") or ""),
                str(st.get("last_daily_reset") or ""),
            ]
        )
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_herentals_ai_context(school_row: Dict[str, Any], max_students: int = 80) -> str:
    """
    Compact read-only snapshot for Gemini grounding (portal school only).
    """
    snap = build_overview_snapshot(school_row)
    students = _load_portal_students(school_row)
    lines = [
        "=== NerdX portal DATA CONTEXT (read-only; do not assume anything not listed) ===",
        f"School name: {school_row.get('name') or ''}",
        f"School code (school_id): {school_row.get('school_id')}",
        f"Portal DB id: {school_row.get('id')}",
        f"Group ID: {school_row.get('group_id') or 'none'}",
        "",
        "Counts and finance (model uses per-learner fee):",
        f"- Portal students: {snap['total_students']}",
        f"- Active last 30 days (daily reset): {snap['active_30d']}",
        f"- Subscription active: {snap['subscription_active']}",
        f"- Subscription days remaining (if active): {snap['days_remaining']}",
        f"- Model monthly revenue USD: {snap['revenue_monthly_model']:.2f}",
        f"- Model school share (30%) USD: {snap['school_share_model']:.2f}",
        f"- Model platform fee (70%) USD: {snap['nerdx_share_model']:.2f}",
        f"- Verified payments total USD: {snap['total_paid']:.2f}",
        f"- Amount due (model fee minus paid) USD: {snap['amount_due']:.2f}",
        "",
    ]
    gid = school_row.get("group_id")
    if gid:
        members = (
            make_supabase_request(
                "GET",
                "schools",
                select="school_id,name,campus_name",
                filters={"group_id": f"eq.{gid}"},
                use_service_role=True,
            )
            or []
        )
        lines.append(f"Group member campuses ({len(members)}):")
        for m in members[:40]:
            lines.append(
                f"  - {m.get('school_id')}: {m.get('name') or ''} ({m.get('campus_name') or 'campus n/a'})"
            )
        if len(members) > 40:
            lines.append(f"  ... and {len(members) - 40} more campuses")
        lines.append("")
    lines.append(f"Students (up to {max_students} of {len(students)}):")
    for st in students[:max_students]:
        nm = f"{st.get('name') or ''} {st.get('surname') or ''}".strip() or "(no name)"
        lines.append(
            f"  - id={st.get('id')}, name={nm!r}, level={st.get('student_level')!r}, "
            f"xp={st.get('xp')}, streak={st.get('streak')}, sub_expires={st.get('subscription_expires_at')}, "
            f"last_daily_reset={st.get('last_daily_reset')}"
        )
    if len(students) > max_students:
        lines.append(f"  ... and {len(students) - max_students} more students not listed")
    return "\n".join(lines)


def generate_board_pptx(school_row: Dict[str, Any]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as e:
        logger.error("python-pptx missing: %s", e)
        raise RuntimeError("PowerPoint export unavailable (python-pptx not installed)") from e

    snap = build_overview_snapshot(school_row)
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = f"{school_row.get('name') or 'School'} — NerdX highlights"
    tf = slide.placeholders[1].text_frame
    tf.text = "Executive snapshot"
    for t in [
        f"Portal students: {snap['total_students']}",
        f"Active in last 30 days: {snap['active_30d']}",
        f"Subscription: {'Active' if snap['subscription_active'] else 'Check renewal'} ({snap['days_remaining']} days left)",
        f"Model school share: ${snap['school_share_model']:.2f} USD / mo",
        f"Verified payments: ${snap['total_paid']:.2f} USD",
    ]:
        p = tf.add_paragraph()
        p.text = t
        p.level = 0
        p.font.size = Pt(18)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
